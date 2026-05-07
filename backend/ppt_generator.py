from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt
from pptx.dml.color import RGBColor

AZUL_ESCURO = RGBColor(19, 52, 87)
BRANCO = RGBColor(255, 255, 255)


# --------------------------------------------------
# SLIDE 1 — BOAS‑VINDAS / INSTITUIÇÃO / DATA
# --------------------------------------------------
def atualizar_slide_1(slide, payload):
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue

        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if "{Instituição}" in run.text:
                    run.text = payload["instituicao"]
                    run.font.name = "Poppins"
                    run.font.size = Pt(20)
                    run.font.bold = True
                    run.font.color.rgb = AZUL_ESCURO
                    paragraph.alignment = PP_ALIGN.CENTER

                if "{Data}" in run.text:
                    run.text = payload["data"]
                    run.font.name = "Poppins"
                    run.font.size = Pt(20)
                    run.font.bold = True
                    run.font.color.rgb = AZUL_ESCURO
                    paragraph.alignment = PP_ALIGN.CENTER


# --------------------------------------------------
# SLIDE 5 — RESUMO DO CONTRATO
# --------------------------------------------------
def atualizar_slide_5(slide, payload):
    resumo = (
        f"Vigência: {payload['vigencia']}\n"
        f"Ambiente: {payload['ambiente']}\n"
        f"Quantidade de licenças: {payload['licencas']}\n"
        f"Catálogos contratados: {payload['catalogos']}\n"
        f"Meses de gestão de licenças: {payload['gestao_licencas']}\n"
        f"Possui aplicativo? {payload['app']}"
    )

    for shape in slide.shapes:
        if shape.has_text_frame and "{RESUMO_CONTRATO}" in shape.text:
            tf = shape.text_frame
            tf.clear()

            for linha in resumo.split("\n"):
                p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                p.text = linha
                for run in p.runs:
                    run.font.name = "Poppins Light"
                    run.font.size = Pt(14)
                    run.font.color.rgb = BRANCO
            break


# --------------------------------------------------
# SLIDE 9 — RESPONSÁVEIS (FONTE 12)
# --------------------------------------------------
def atualizar_slide_9(slide, payload):
    linhas = [
        "Comercial Responsável:",
        payload["comercial"]["nome"],
        f"E-mail: {payload['comercial']['email']}",
        f"Telefone: {payload['comercial']['telefone']}",
        "",
        "Customer Success Responsável:",
        payload["cs"]["nome"],
        f"E-mail: {payload['cs']['email']}",
        f"Telefone: {payload['cs']['telefone']}",
    ]

    for shape in slide.shapes:
        if shape.has_text_frame and "{RESPONSAVEIS}" in shape.text:
            tf = shape.text_frame
            tf.clear()

            for linha in linhas:
                p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                p.text = linha
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Poppins Light"
                    run.font.size = Pt(12)   # ✅ AUMENTADO PARA 12
                    run.font.color.rgb = AZUL_ESCURO
            break


# --------------------------------------------------
# SLIDE 17 — EMAILS DE ACESSO (Poppins Light 9)
# --------------------------------------------------
def atualizar_slide_17(slide, payload):
    emails = payload.get("emails_acesso", [])

    for shape in slide.shapes:
        if shape.has_text_frame and "{EMAILS_ACESSO}" in shape.text:
            tf = shape.text_frame
            tf.clear()

            for email in emails:
                p = tf.add_paragraph() if tf.text else tf.paragraphs[0]
                p.text = email
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = "Poppins Light"  # ✅ Light
                    run.font.size = Pt(9)           # ✅ tamanho 9
            break


# --------------------------------------------------
# FUNÇÃO PRINCIPAL
# --------------------------------------------------
def generate_ppt(template_path, output_path, payload):
    prs = Presentation(template_path)

    atualizar_slide_1(prs.slides[0], payload)
    atualizar_slide_5(prs.slides[4], payload)
    atualizar_slide_9(prs.slides[8], payload)
    atualizar_slide_17(prs.slides[16], payload)

    prs.save(output_path)